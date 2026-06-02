"""
generate_logic_ppt.py
MR FTE算出システム ロジック説明 PowerPoint生成スクリプト
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree


# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
C_NAVY      = RGBColor(0x1a, 0x2a, 0x4a)   # dark navy – primary / bg
C_LIGHT_BG  = RGBColor(0xee, 0xf2, 0xf9)   # light navy – slide bg
C_BLUE      = RGBColor(0x25, 0x63, 0xeb)   # secondary blue
C_GREEN     = RGBColor(0x10, 0xb9, 0x81)   # accent green
C_AMBER     = RGBColor(0xf5, 0x9e, 0x0b)   # warning amber
C_WHITE     = RGBColor(0xff, 0xff, 0xff)
C_DARK_TEXT = RGBColor(0x1e, 0x29, 0x3b)
C_MID_GRAY  = RGBColor(0x64, 0x74, 0x8b)
C_LIGHT_BLUE_BOX = RGBColor(0xdb, 0xe9, 0xff)
C_GREEN_BOX = RGBColor(0xd1, 0xfa, 0xe5)
C_AMBER_BOX = RGBColor(0xfe, 0xf3, 0xc7)

# Slide size: 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

FONT_JP = "Meiryo"
FONT_EN = "Calibri"


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs: Presentation):
    return prs.slide_layouts[6]   # completely blank


def rgb(r, g, b) -> RGBColor:
    return RGBColor(r, g, b)


def fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line=False):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill_solid(shape, fill_color)
    if not line:
        no_line(shape)
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=C_DARK_TEXT,
                align=PP_ALIGN.LEFT, font=FONT_JP, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size     = Pt(font_size)
    run.font.bold     = bold
    run.font.color.rgb = color
    run.font.name     = font
    return txb


def add_label_textbox(slide, text, left, top, width, height,
                      font_size=14, bold=False,
                      text_color=C_WHITE, bg_color=C_NAVY):
    """Colored background box with text."""
    rect = add_rect(slide, left, top, width, height, bg_color)
    tf   = rect.text_frame
    tf.word_wrap = True
    p    = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run  = p.add_run()
    run.text = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = text_color
    run.font.name      = FONT_JP
    # vertical centering
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return rect


def set_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text, is_dark_slide=False):
    """Add narrow navy bar at top with title text."""
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.45), C_NAVY)
    tf  = bar.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + title_text
    run.font.size      = Pt(15)
    run.font.bold      = True
    run.font.color.rgb = C_WHITE
    run.font.name      = FONT_JP
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_slide_number(slide, number: int):
    """Add slide number at bottom-right."""
    add_textbox(
        slide, str(number),
        left=Inches(12.6), top=Inches(7.1),
        width=Inches(0.6), height=Inches(0.3),
        font_size=11, color=C_MID_GRAY,
        align=PP_ALIGN.RIGHT, font=FONT_EN
    )


def add_section_title(slide, text, top=Inches(0.55)):
    """Big section heading below header bar."""
    add_textbox(
        slide, text,
        left=Inches(0.4), top=top,
        width=Inches(12.5), height=Inches(0.55),
        font_size=22, bold=True, color=C_NAVY,
        align=PP_ALIGN.LEFT
    )


# ===========================================================================
# SLIDE 1 – Title
# ===========================================================================
def slide1_title(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_NAVY)

    # Decorative accent bar (left side)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, C_BLUE)

    # Green accent stripe
    add_rect(slide, Inches(0.18), Inches(3.6), Inches(9), Inches(0.06), C_GREEN)

    # Main title
    t = add_textbox(
        slide, "MR FTE算出システム\nロジック説明",
        left=Inches(0.55), top=Inches(1.4),
        width=Inches(11.5), height=Inches(1.8),
        font_size=40, bold=True, color=C_WHITE,
        align=PP_ALIGN.LEFT
    )
    t.text_frame.paragraphs[0].runs[0].font.name = FONT_JP

    # Subtitle
    add_textbox(
        slide, "FY2026-2035  リソースアロケーション シミュレーション",
        left=Inches(0.55), top=Inches(3.75),
        width=Inches(11.5), height=Inches(0.5),
        font_size=20, color=rgb(0xb8, 0xce, 0xf9),
        align=PP_ALIGN.LEFT
    )

    # Bottom note box
    note_box = add_rect(
        slide,
        Inches(0.55), Inches(6.3), Inches(6.5), Inches(0.6),
        rgb(0x0f, 0x1c, 0x36)
    )
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  CS: 380 FTE / PS: 45 FTE  制約下の最適配分"
    run.font.size      = Pt(13)
    run.font.color.rgb = C_AMBER
    run.font.bold      = True
    run.font.name      = FONT_JP
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_slide_number(slide, 1)


# ===========================================================================
# SLIDE 2 – 全体フロー概要
# ===========================================================================
def slide2_flow(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "MR FTE算出システム  ロジック説明")
    add_section_title(slide, "算出ステップ 全体フロー")

    steps = [
        ("STEP 1", "ターゲット医師数",   "R/W ティア CSV 直接管理",           C_BLUE),
        ("STEP 2", "FC/SC 分割",        "品目重複・品目セット考慮",            C_BLUE),
        ("STEP 3", "訪問頻度",           "ライフサイクル調整",                 C_BLUE),
        ("STEP 4", "FTE算出",            "required_calls ÷ 月次キャパシティ",  C_BLUE),
        ("STEP 5", "HC 正規化",          "エリア×月の制約 (380/45)",           C_BLUE),
        ("STEP 6", "限界 ROI 配分",      "新製品への最適 FTE 配分",            C_GREEN),
        ("別軸",   "デジタル有効性スコア","0–1 スコアでチャネル戦略示唆",       C_AMBER),
    ]

    box_w   = Inches(1.55)
    box_h   = Inches(0.78)
    arr_w   = Inches(0.18)
    gap     = Inches(0.12)
    total_w = len(steps) * box_w + (len(steps) - 1) * (arr_w + gap * 2)
    start_x = (SLIDE_W - total_w) / 2
    y_main  = Inches(1.65)
    y_sub   = Inches(2.55)

    for i, (label, name, detail, color) in enumerate(steps):
        x = start_x + i * (box_w + arr_w + gap * 2)

        # Step box
        box = add_rect(slide, x, y_main, box_w, box_h, color)
        tf  = box.text_frame
        tf.word_wrap = True
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run()
        r1.text = label + "\n"
        r1.font.size      = Pt(9)
        r1.font.bold      = True
        r1.font.color.rgb = C_WHITE
        r1.font.name      = FONT_JP
        r2 = p.add_run()
        r2.text = name
        r2.font.size      = Pt(11)
        r2.font.bold      = True
        r2.font.color.rgb = C_WHITE
        r2.font.name      = FONT_JP

        # Detail text below
        add_textbox(
            slide, detail,
            left=x, top=y_sub,
            width=box_w, height=Inches(0.55),
            font_size=9, color=C_DARK_TEXT,
            align=PP_ALIGN.CENTER
        )

        # Arrow (except last)
        if i < len(steps) - 1:
            ax = x + box_w + gap
            ay = y_main + box_h / 2 - Inches(0.05)
            add_rect(slide, ax, ay, arr_w, Inches(0.1),
                     C_MID_GRAY if color != C_AMBER else C_AMBER)

    # Add a dotted separator before the "別軸" item
    sep_x = start_x + 6 * (box_w + arr_w + gap * 2) - gap
    add_rect(slide, sep_x - Inches(0.02), y_main - Inches(0.15),
             Inches(0.04), box_h + Inches(0.3), C_AMBER)

    # Note at bottom
    add_textbox(
        slide,
        "※ 別軸のデジタル有効性スコアは MR FTE とは独立して算出。FTE コスト計上なし。",
        left=Inches(0.5), top=Inches(3.55),
        width=Inches(12.3), height=Inches(0.4),
        font_size=11, color=C_MID_GRAY,
        align=PP_ALIGN.LEFT
    )

    add_slide_number(slide, 2)


# ===========================================================================
# SLIDE 3 – STEP 1 ターゲット医師数
# ===========================================================================
def slide3_step1(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "MR FTE算出システム  ロジック説明")
    add_section_title(slide, "STEP 1: ターゲット医師数の決定")

    # Three tier boxes
    tiers = [
        ("R 医師",   "患者数かなり多め",  "訪問頻度  2.0 回/月", C_BLUE,  C_LIGHT_BLUE_BOX),
        ("W 医師",   "患者数多め",        "訪問頻度  1.0 回/月", C_GREEN, C_GREEN_BOX),
    ]
    bw, bh = Inches(3.4), Inches(1.1)
    for i, (tier, desc, freq, hdr_c, bg_c) in enumerate(tiers):
        x = Inches(0.5) + i * (bw + Inches(0.35))
        y = Inches(1.35)
        box = add_rect(slide, x, y, bw, bh, bg_c)
        no_line(box)
        # header strip
        add_rect(slide, x, y, bw, Inches(0.3), hdr_c)
        # header text
        add_label_textbox(slide, tier, x, y, bw, Inches(0.3),
                          font_size=12, bold=True, text_color=C_WHITE, bg_color=hdr_c)
        # body text
        add_textbox(slide, desc + "\n" + freq,
                    x + Inches(0.1), y + Inches(0.35),
                    bw - Inches(0.2), Inches(0.7),
                    font_size=13, color=C_DARK_TEXT)

    # Formula box
    formula_box = add_rect(slide, Inches(7.8), Inches(1.35),
                           Inches(5.0), Inches(1.1), rgb(0x1e, 0x40, 0xaf))
    tf = formula_box.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "target_drs(p) =\nr_doctors(p)  +  w_doctors(p)"
    r.font.size      = Pt(14)
    r.font.bold      = True
    r.font.color.rgb = C_WHITE
    r.font.name      = FONT_JP

    # Yearly note
    note = add_rect(slide, Inches(0.5), Inches(2.65), Inches(12.3), Inches(0.55),
                    rgb(0xff, 0xf7, 0xe0))
    no_line(note)
    tf = note.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "  年度別変化は target_doctor_yearly.csv で指定（新製品ランプアップ・LOE 後縮小）"
    r.font.size      = Pt(12)
    r.font.color.rgb = rgb(0x92, 0x40, 0x00)
    r.font.name      = FONT_JP

    # Example table
    add_textbox(slide, "【例】FY2026/4 基準値",
                Inches(0.5), Inches(3.35),
                Inches(5), Inches(0.35),
                font_size=13, bold=True, color=C_NAVY)

    headers = ["品目", "R 医師", "W 医師", "合計"]
    rows    = [
        ["GLI",  "1,200", "3,000", "4,200"],
        ["OVE",  "  280", "  820", "1,100"],
        ["LVM",  "  200", "  600", "  800"],
    ]
    col_w  = [Inches(1.5), Inches(1.2), Inches(1.2), Inches(1.2)]
    col_xs = [Inches(0.5), Inches(2.0), Inches(3.2), Inches(4.4)]
    row_h  = Inches(0.38)
    row_ys = [Inches(3.72), Inches(4.1), Inches(4.48), Inches(4.86)]

    # Header row
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_w)):
        cell = add_rect(slide, cx, row_ys[0], cw, row_h, C_NAVY)
        tf   = cell.text_frame
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hdr
        r.font.size      = Pt(12)
        r.font.bold      = True
        r.font.color.rgb = C_WHITE
        r.font.name      = FONT_JP

    # Data rows
    for ri, row in enumerate(rows):
        bg = C_WHITE if ri % 2 == 0 else rgb(0xf1, 0xf5, 0xf9)
        for j, (val, cx, cw) in enumerate(zip(row, col_xs, col_w)):
            cell = add_rect(slide, cx, row_ys[ri + 1], cw, row_h, bg)
            tf   = cell.text_frame
            from pptx.enum.text import MSO_ANCHOR
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = ("  " if j == 0 else "") + val
            r.font.size      = Pt(12)
            r.font.color.rgb = C_DARK_TEXT
            r.font.name      = FONT_JP

    add_slide_number(slide, 3)


# ===========================================================================
# SLIDE 4 – STEP 2 & 3 FC/SC と訪問頻度
# ===========================================================================
def slide4_step23(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "MR FTE算出システム  ロジック説明")
    add_section_title(slide, "STEP 2-3: FC/SC 分割と訪問頻度")

    col_y  = Inches(1.3)
    col_h  = Inches(5.5)
    col_w  = Inches(5.9)
    gap    = Inches(0.4)
    left_x = Inches(0.4)
    right_x= left_x + col_w + gap

    # Left column bg
    lbg = add_rect(slide, left_x, col_y, col_w, col_h, C_WHITE)
    no_line(lbg)
    # Right column bg
    rbg = add_rect(slide, right_x, col_y, col_w, col_h, C_WHITE)
    no_line(rbg)

    # Left header
    add_label_textbox(slide, "FC / SC 分割",
                      left_x, col_y, col_w, Inches(0.38),
                      font_size=14, bold=True,
                      text_color=C_WHITE, bg_color=C_BLUE)

    # Right header
    add_label_textbox(slide, "訪問頻度の決定",
                      right_x, col_y, col_w, Inches(0.38),
                      font_size=14, bold=True,
                      text_color=C_WHITE, bg_color=C_GREEN)

    fc_items = [
        ("FC ドクター", "自品目メイン訪問 → フルコスト計上", C_BLUE),
        ("SC ドクター", "他品目メインのついで訪問\n→ × SC 係数 (0.1)", C_MID_GRAY),
        ("品目セット判定", "activity_set.csv で重複チェック\n被り品目は SC 扱いに自動振替", C_NAVY),
    ]
    for i, (ttl, body, color) in enumerate(fc_items):
        y = col_y + Inches(0.55) + i * Inches(1.55)
        add_rect(slide, left_x + Inches(0.1), y,
                 Inches(0.07), Inches(0.9), color)
        add_textbox(slide, ttl,
                    left_x + Inches(0.28), y,
                    col_w - Inches(0.38), Inches(0.35),
                    font_size=13, bold=True, color=color)
        add_textbox(slide, body,
                    left_x + Inches(0.28), y + Inches(0.35),
                    col_w - Inches(0.38), Inches(0.6),
                    font_size=11, color=C_DARK_TEXT)

    freq_items = [
        ("実効頻度の算式",
         "実効頻度 = tier_freq × lifecycle_adj",
         C_GREEN),
        ("LOE 後 小分子",
         "lifecycle_adj = 0.0\n（活動終了・FTE ゼロ化）",
         C_AMBER),
        ("LOE 後 バイオ",
         "lifecycle_adj = 0.55\n（バイオシミラー対策で継続活動）",
         C_BLUE),
    ]
    for i, (ttl, body, color) in enumerate(freq_items):
        y = col_y + Inches(0.55) + i * Inches(1.55)
        add_rect(slide, right_x + Inches(0.1), y,
                 Inches(0.07), Inches(0.9), color)
        add_textbox(slide, ttl,
                    right_x + Inches(0.28), y,
                    col_w - Inches(0.38), Inches(0.35),
                    font_size=13, bold=True, color=color)
        add_textbox(slide, body,
                    right_x + Inches(0.28), y + Inches(0.35),
                    col_w - Inches(0.38), Inches(0.6),
                    font_size=11, color=C_DARK_TEXT)

    add_slide_number(slide, 4)


# ===========================================================================
# SLIDE 5 – STEP 4 FTE算出
# ===========================================================================
def slide5_step4(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "MR FTE算出システム  ロジック説明")
    add_section_title(slide, "STEP 4: 必要 FTE 算出")

    # Main formula block
    formula_bg = add_rect(slide, Inches(0.5), Inches(1.3),
                          Inches(12.3), Inches(2.1), C_NAVY)
    tf = formula_bg.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    lines = [
        "  required_calls  =  R医師 × r_freq × lc_adj  +  W医師 × w_freq × lc_adj",
        "  base_fte            =  required_calls  ÷  (calls_per_day × working_days/month)",
    ]
    for li, line in enumerate(lines):
        if li == 0:
            r = p.add_run()
        else:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            r = p2.add_run()
        r.text = line
        r.font.size      = Pt(14)
        r.font.color.rgb = C_GREEN if li == 0 else rgb(0xb8, 0xce, 0xf9)
        r.font.bold      = True
        r.font.name      = FONT_JP

    # CS / PS capacity boxes
    cap_boxes = [
        ("CS キャパシティ",
         "2.5 calls/day × 20 days\n= 50 calls/month/FTE", C_BLUE),
        ("PS キャパシティ",
         "1.5 calls/day × 20 days\n= 30 calls/month/FTE", C_GREEN),
    ]
    for i, (ttl, body, color) in enumerate(cap_boxes):
        x = Inches(0.5) + i * Inches(4.2)
        box = add_rect(slide, x, Inches(3.55), Inches(3.8), Inches(1.0), color)
        tf = box.text_frame
        tf.word_wrap = True
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run()
        r1.text = ttl + "\n"
        r1.font.size      = Pt(12)
        r1.font.bold      = True
        r1.font.color.rgb = C_WHITE
        r1.font.name      = FONT_JP
        r2 = p.add_run()
        r2.text = body
        r2.font.size      = Pt(13)
        r2.font.bold      = True
        r2.font.color.rgb = C_WHITE
        r2.font.name      = FONT_JP

    # Boost factors
    add_textbox(slide, "× ブーストファクター",
                Inches(0.5), Inches(4.75),
                Inches(5), Inches(0.38),
                font_size=14, bold=True, color=C_AMBER)

    boosts = [
        "効能追加ブースト : +30% を最大 12 ヶ月間適用",
        "競合品ブースト   : 競合発売直後に活動強化",
        "供給制限ファクター : 供給不足時は FTE を削減",
    ]
    for i, b in enumerate(boosts):
        y = Inches(5.2) + i * Inches(0.42)
        add_rect(slide, Inches(0.5), y + Inches(0.08),
                 Inches(0.12), Inches(0.22), C_AMBER)
        add_textbox(slide, b,
                    Inches(0.75), y,
                    Inches(11.5), Inches(0.38),
                    font_size=13, color=C_DARK_TEXT)

    add_slide_number(slide, 5)


# ===========================================================================
# SLIDE 6 – STEP 5 & 6 HC正規化と限界ROI
# ===========================================================================
def slide6_step56(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "MR FTE算出システム  ロジック説明")
    add_section_title(slide, "STEP 5-6: HC 正規化と新製品 FTE 配分")

    # ---- STEP 5 HC正規化 ----
    s5_x, s5_y = Inches(0.4), Inches(1.3)
    s5_w, s5_h = Inches(5.9), Inches(5.5)
    bg5 = add_rect(slide, s5_x, s5_y, s5_w, s5_h, C_WHITE)
    no_line(bg5)
    add_label_textbox(slide, "STEP 5  HC 正規化",
                      s5_x, s5_y, s5_w, Inches(0.38),
                      font_size=14, bold=True,
                      text_color=C_WHITE, bg_color=C_BLUE)

    s5_items = [
        "エリア総 FTE を制約値に正規化",
        "CS:  Σ required_fte  →  380 FTE に正規化",
        "PS:  Σ required_fte  →   45 FTE に正規化",
    ]
    for i, txt in enumerate(s5_items):
        add_textbox(slide, ("• " if i == 0 else "  ") + txt,
                    s5_x + Inches(0.2), s5_y + Inches(0.55) + i * Inches(0.52),
                    s5_w - Inches(0.4), Inches(0.45),
                    font_size=12, color=C_DARK_TEXT,
                    bold=(i == 0))

    # Formula
    f5 = add_rect(slide, s5_x + Inches(0.15), s5_y + Inches(2.2),
                  s5_w - Inches(0.3), Inches(0.85), C_NAVY)
    tf = f5.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "adjusted_fte(p) =\nrequired_fte(p) × (constraint / Σ required_fte)"
    r.font.size      = Pt(12)
    r.font.bold      = True
    r.font.color.rgb = C_GREEN
    r.font.name      = FONT_JP

    # ---- STEP 6 限界ROI ----
    s6_x = s5_x + s5_w + Inches(0.4)
    bg6 = add_rect(slide, s6_x, s5_y, s5_w, s5_h, C_WHITE)
    no_line(bg6)
    add_label_textbox(slide, "STEP 6  限界 ROI 最適配分",
                      s6_x, s5_y, s5_w, Inches(0.38),
                      font_size=14, bold=True,
                      text_color=C_WHITE, bg_color=C_GREEN)

    s6_items = [
        "新製品 FTE をドナー品目から最適配分",
        "Hill 関数でレスポンスカーブを近似",
        "dSales/dFTE が等しくなるまで最適化",
        "ドナー: FY2026/4 基準 FTE の 50% まで削減可",
    ]
    for i, txt in enumerate(s6_items):
        add_textbox(slide, "• " + txt,
                    s6_x + Inches(0.2), s5_y + Inches(0.55) + i * Inches(0.52),
                    s5_w - Inches(0.4), Inches(0.45),
                    font_size=12, color=C_DARK_TEXT)

    # Hill formula
    f6 = add_rect(slide, s6_x + Inches(0.15), s5_y + Inches(2.7),
                  s5_w - Inches(0.3), Inches(1.0), C_NAVY)
    tf = f6.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "sales_lift =\nk × FTEⁿ / (EC50ⁿ + FTEⁿ)"
    r.font.size      = Pt(13)
    r.font.bold      = True
    r.font.color.rgb = C_AMBER
    r.font.name      = FONT_JP

    add_slide_number(slide, 6)


# ===========================================================================
# SLIDE 7 – デジタル有効性スコア
# ===========================================================================
def slide7_digital(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "MR FTE算出システム  ロジック説明")
    add_section_title(slide, "デジタル有効性スコア（MR FTE とは独立）")

    # Independence note
    note = add_rect(slide, Inches(0.5), Inches(1.3),
                    Inches(12.3), Inches(0.55), rgb(0xfe, 0xf3, 0xc7))
    no_line(note)
    tf = note.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "  デジタルは FTE コストなし（m3 等への掲載）→ MR FTE と完全分離して算出"
    r.font.size      = Pt(13)
    r.font.color.rgb = rgb(0x92, 0x40, 0x00)
    r.font.bold      = True
    r.font.name      = FONT_JP

    # Formula
    f_box = add_rect(slide, Inches(0.5), Inches(2.0),
                     Inches(12.3), Inches(0.75), C_NAVY)
    tf = f_box.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "score = 0.5 × MMM_digital_fraction  +  0.5 × digital_soc_rate  +  lifecycle_adj"
    r.font.size      = Pt(14)
    r.font.bold      = True
    r.font.color.rgb = C_GREEN
    r.font.name      = FONT_JP

    # Level thresholds
    add_textbox(slide, "スコアに基づくチャネル戦略",
                Inches(0.5), Inches(2.95),
                Inches(7), Inches(0.38),
                font_size=13, bold=True, color=C_NAVY)

    thresholds = [
        ("高 (≥ 0.55)", "デジタルチャネルを積極活用",   C_GREEN, rgb(0xd1, 0xfa, 0xe5)),
        ("中 (≥ 0.35)", "MR とデジタルのバランス戦略",  C_BLUE,  rgb(0xdb, 0xe9, 0xff)),
        ("低 (< 0.35)", "MR 中心戦略を推奨",            C_AMBER, rgb(0xfe, 0xf3, 0xc7)),
    ]
    for i, (lvl, desc, lbl_c, bg_c) in enumerate(thresholds):
        y = Inches(3.4) + i * Inches(0.62)
        # Label
        lbl = add_rect(slide, Inches(0.5), y, Inches(1.6), Inches(0.48), lbl_c)
        tf = lbl.text_frame
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = lvl
        r.font.size      = Pt(12)
        r.font.bold      = True
        r.font.color.rgb = C_WHITE
        r.font.name      = FONT_JP
        # Description
        desc_box = add_rect(slide, Inches(2.2), y, Inches(4.8), Inches(0.48), bg_c)
        no_line(desc_box)
        tf = desc_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "  " + desc
        r.font.size      = Pt(12)
        r.font.color.rgb = C_DARK_TEXT
        r.font.name      = FONT_JP

    # Lifecycle adjustments
    add_textbox(slide, "ライフサイクル調整（lifecycle_adj）",
                Inches(7.5), Inches(2.95),
                Inches(5.3), Inches(0.38),
                font_size=13, bold=True, color=C_NAVY)

    lc_items = [
        ("LOE 後",       "− 0.25  デジタル費用対効果低下", C_AMBER),
        ("成長期 1-3年", "+ 0.08  新薬認知向上に有効",      C_GREEN),
    ]
    for i, (phase, adj, color) in enumerate(lc_items):
        y = Inches(3.4) + i * Inches(0.7)
        add_rect(slide, Inches(7.5), y + Inches(0.08),
                 Inches(0.1), Inches(0.32), color)
        add_textbox(slide, phase,
                    Inches(7.7), y,
                    Inches(2.2), Inches(0.35),
                    font_size=12, bold=True, color=color)
        add_textbox(slide, adj,
                    Inches(7.7), y + Inches(0.33),
                    Inches(5.1), Inches(0.35),
                    font_size=11, color=C_DARK_TEXT)

    add_slide_number(slide, 7)


# ===========================================================================
# SLIDE 8 – データファイル一覧
# ===========================================================================
def slide8_files(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_LIGHT_BG)
    add_header_bar(slide, "MR FTE算出システム  ロジック説明")
    add_section_title(slide, "主要入力データファイル")

    headers = ["ファイル名", "内容", "更新頻度"]
    rows = [
        ["target_doctors.csv",        "R/W 医師数・訪問頻度",      "定期見直し"],
        ["target_doctor_yearly.csv",  "年度別医師数推移",           "定期見直し"],
        ["products.csv",              "品目設定（発売日・LOE 等）", "変更時"],
        ["mmm_decay_params.csv",      "MMM パラメータ",             "年次更新"],
        ["soc_params.csv",            "SOC 想起率",                 "年次更新"],
        ["current_fte.csv",           "FY2026/4 基準 FTE",          "期初確認"],
        ["sales_forecast.csv",        "売上予測",                   "四半期更新"],
        ["competitor_schedule.csv",   "競合品発売予定",             "随時更新"],
    ]

    col_ws = [Inches(3.8), Inches(5.6), Inches(2.8)]
    col_xs = [Inches(0.45),
              Inches(0.45) + col_ws[0],
              Inches(0.45) + col_ws[0] + col_ws[1]]
    row_h = Inches(0.5)
    row_ys = [Inches(1.35) + i * row_h for i in range(len(rows) + 1)]

    # Header row
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
        cell = add_rect(slide, cx, row_ys[0], cw, row_h, C_NAVY)
        tf   = cell.text_frame
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hdr
        r.font.size      = Pt(13)
        r.font.bold      = True
        r.font.color.rgb = C_WHITE
        r.font.name      = FONT_JP

    # Data rows
    for ri, row in enumerate(rows):
        bg = C_WHITE if ri % 2 == 0 else rgb(0xf1, 0xf5, 0xf9)
        for j, (val, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
            cell = add_rect(slide, cx, row_ys[ri + 1], cw, row_h, bg)
            tf   = cell.text_frame
            from pptx.enum.text import MSO_ANCHOR
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "  " + val
            r.font.size = Pt(11 if j == 0 else 12)
            r.font.color.rgb = C_BLUE if j == 0 else C_DARK_TEXT
            r.font.bold      = (j == 0)
            r.font.name      = FONT_JP

    add_slide_number(slide, 8)


# ===========================================================================
# SLIDE 9 – サマリー・主要アウトプット
# ===========================================================================
def slide9_summary(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_slide_bg(slide, C_NAVY)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.45), rgb(0x0f, 0x1c, 0x36))

    add_textbox(slide, "  シミュレーション出力サマリー",
                Inches(0), Inches(0),
                Inches(13), Inches(0.45),
                font_size=15, bold=True, color=C_WHITE,
                align=PP_ALIGN.LEFT)

    # Big stats row
    stats = [
        ("19 品目",           "FY2026-2035 対象品目数",  C_GREEN),
        ("10 年間 / 120 ヶ月","シミュレーション期間",     C_BLUE),
        ("CS 380 / PS 45",    "制約 FTE (最適配分)",      C_AMBER),
    ]
    sw = Inches(3.9)
    sx = Inches(0.5)
    sy = Inches(0.7)
    for i, (val, lbl, color) in enumerate(stats):
        x = sx + i * (sw + Inches(0.15))
        box = add_rect(slide, x, sy, sw, Inches(1.2), rgb(0x0f, 0x1c, 0x36))
        no_line(box)
        add_rect(slide, x, sy, Inches(0.15), Inches(1.2), color)
        add_textbox(slide, val,
                    x + Inches(0.25), sy + Inches(0.05),
                    sw - Inches(0.3), Inches(0.65),
                    font_size=22, bold=True, color=color,
                    align=PP_ALIGN.LEFT)
        add_textbox(slide, lbl,
                    x + Inches(0.25), sy + Inches(0.65),
                    sw - Inches(0.3), Inches(0.45),
                    font_size=12, color=rgb(0x94, 0xa3, 0xb8),
                    align=PP_ALIGN.LEFT)

    # HTML outputs
    add_textbox(slide, "主要 HTML 出力ファイル",
                Inches(0.5), Inches(2.1),
                Inches(10), Inches(0.45),
                font_size=15, bold=True, color=C_WHITE)

    html_files = [
        ("fy2029_fte_report.html",    "品目別 FTE 推移・デジタル有効性スコア一覧",       C_GREEN),
        ("fte_logic_document.html",   "本ロジック詳細ドキュメント（HTML 版）",            C_BLUE),
        ("fy2029_sensitivity.html",   "感度分析（±20% シナリオ・主要品目影響度比較）",   C_AMBER),
    ]
    for i, (fname, desc, color) in enumerate(html_files):
        y = Inches(2.65) + i * Inches(1.3)
        box = add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.1),
                       rgb(0x0f, 0x1c, 0x36))
        no_line(box)
        add_rect(slide, Inches(0.5), y, Inches(0.12), Inches(1.1), color)
        add_textbox(slide, fname,
                    Inches(0.75), y + Inches(0.1),
                    Inches(11.8), Inches(0.42),
                    font_size=14, bold=True, color=color,
                    align=PP_ALIGN.LEFT)
        add_textbox(slide, desc,
                    Inches(0.75), y + Inches(0.55),
                    Inches(11.8), Inches(0.42),
                    font_size=12, color=rgb(0xb8, 0xce, 0xf9),
                    align=PP_ALIGN.LEFT)

    add_slide_number(slide, 9)


# ===========================================================================
# MAIN
# ===========================================================================
def main():
    out_dir = Path(__file__).parent.parent / "data" / "output"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "fte_logic.pptx"

    prs = new_prs()

    slide1_title(prs)
    slide2_flow(prs)
    slide3_step1(prs)
    slide4_step23(prs)
    slide5_step4(prs)
    slide6_step56(prs)
    slide7_digital(prs)
    slide8_files(prs)
    slide9_summary(prs)

    prs.save(str(out_path))
    print(f"出力完了: {out_path}")


if __name__ == "__main__":
    main()
